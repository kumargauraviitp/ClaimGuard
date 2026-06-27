"""
ClaimGuard — research presentation builder.
Produces ClaimGuard_Research.pptx (11 slides) for a research-paper-author
audience.  Every figure / number is drawn from the project's real artifacts
(features.json, leaderboard.csv, training_report.md, dataset_analysis_report.md,
serving.py, engine.py, graph.py).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "ppt_assets")
OUT = os.path.join(HERE, "ClaimGuard_Research.pptx")

# ---------- palette ----------
NAVY  = RGBColor(0x0B, 0x1F, 0x3A)
TEAL  = RGBColor(0x0E, 0x7C, 0x86)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
SLATE = RGBColor(0x5D, 0x6D, 0x7E)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x21, 0x29, 0x33)

# 16:9
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def band(slide, top, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SW, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
         space_after=6, line_spacing=1.0):
    """runs: str OR list of (text, opts) tuples; opts dict overrides."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(runs, str):
        runs = [(runs, {})]
    lines = []
    for item in runs:
        if isinstance(item, str):
            lines.append((item, {}))
        else:
            lines.append(item)
    first = True
    for content, opts in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space_after", space_after))
        p.space_before = Pt(opts.get("space_before", 0))
        p.line_spacing = opts.get("line_spacing", line_spacing)
        # content may itself be a list of inline runs
        if isinstance(content, str):
            content = [content]
        for seg in content:
            if isinstance(seg, tuple):
                seg_text, seg_opts = seg
            else:
                seg_text, seg_opts = seg, {}
            r = p.add_run(); r.text = seg_text
            r.font.size = Pt(opts.get("size", seg_opts.get("size", size)))
            r.font.bold = opts.get("bold", seg_opts.get("bold", bold))
            r.font.name = seg_opts.get("font", font)
            r.font.color.rgb = opts.get("color", seg_opts.get("color", color))
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=INK, bullet="▪",
            bcolor=TEAL, gap=8, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    first = True
    for item in items:
        # item may be (text, level) or (text, level, opts)
        level = 0; opts = {}
        if isinstance(item, tuple):
            if len(item) == 2:
                txt, level = item
            else:
                txt, level, opts = item
        else:
            txt = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(opts.get("space_after", gap))
        p.space_before = Pt(opts.get("space_before", 0))
        p.level = level
        rb = p.add_run(); rb.text = (("    " * level) +
                                     ("" if level else bullet + "  "))
        rb.font.size = Pt(opts.get("size", size))
        rb.font.color.rgb = opts.get("bcolor", bcolor if level == 0 else SLATE)
        rb.font.bold = True
        rt = p.add_run(); rt.text = txt
        rt.font.size = Pt(opts.get("size", size))
        rt.font.color.rgb = opts.get("color", color)
        rt.font.bold = opts.get("bold", False)
    return tb


def card(slide, x, y, w, h, fill=LIGHT, line=SLATE, line_w=0.75, radius=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = slide.shapes.add_shape(shp_type, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = line; c.line.width = Pt(line_w)
    c.shadow.inherit = False
    if radius:
        try:
            c.adjustments[0] = 0.06
        except Exception:
            pass
    return c


def kpi(slide, x, y, w, h, value, label, vcolor=NAVY, fill=LIGHT):
    card(slide, x, y, w, h, fill=fill, line=RGBColor(0xDD,0xE3,0xEA))
    text(slide, x, y+Emu(int(h*0.16)), w, Inches(0.6), value,
         size=30, bold=True, color=vcolor, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x, y+Emu(int(h*0.66)), w, Inches(0.5), label,
         size=12, color=SLATE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def pic(slide, path, x, y, w=None, h=None):
    kw = {}
    if w is not None: kw["width"] = w
    if h is not None: kw["height"] = h
    return slide.shapes.add_picture(path, x, y, **kw)


def footer(slide, n):
    band(slide, Inches(7.2), Inches(0.3), WHITE)
    text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.3),
         "ClaimGuard — AI-Powered Insurance Fraud Detection", size=9,
         color=SLATE)
    text(slide, Inches(12.0), Inches(7.12), Inches(0.9), Inches(0.3),
         str(n), size=9, color=SLATE, align=PP_ALIGN.RIGHT)


def header(slide, kicker, title):
    band(slide, 0, Inches(1.5), NAVY)  # thin top accent via shape below
    # kicker
    text(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.3),
         kicker.upper(), size=11, bold=True, color=TEAL)
    # title
    text(slide, Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.7),
         title, size=27, bold=True, color=NAVY)
    # rule under title
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.5), Inches(1.34),
                                Inches(2.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = AMBER
    ln.line.fill.background(); ln.shadow.inherit = False


# ==================================================================
# SLIDE 1 — TITLE
# ==================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# decorative side block
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.7), 0,
                       Inches(3.633), SH)
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x10, 0x2A, 0x4D)
b.line.fill.background(); b.shadow.inherit = False
# accent bar
ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6),
                        Inches(1.2), Pt(5))
ab.fill.solid(); ab.fill.fore_color.rgb = AMBER
ab.line.fill.background(); ab.shadow.inherit = False
text(s, Inches(0.8), Inches(1.7), Inches(8.5), Inches(0.5),
     "RESEARCH PRESENTATION", size=13, bold=True, color=TEAL)
text(s, Inches(0.8), Inches(2.8), Inches(9.0), Inches(1.8),
     "ClaimGuard", size=58, bold=True, color=WHITE)
text(s, Inches(0.8), Inches(3.9), Inches(9.0), Inches(1.2),
     "Explainable, Multi-Agent Machine Learning for\n"
     "Real-Time Insurance Fraud Detection",
     size=24, color=RGBColor(0xCF,0xDD,0xEA))
# authors / meta block
text(s, Inches(0.8), Inches(5.5), Inches(8.5), Inches(1.4),
     [("System Design, Modeling & Deployment", {"size":15,"color":WHITE,
        "bold":True,"space_after":4}),
      ("Full-stack platform: FastAPI · Next.js · LightGBM · SHAP · "
       "LangGraph", {"size":12,"color":RGBColor(0x9F,0xB6,0xCC)}),
      ("June 2026", {"size":12,"color":RGBColor(0x9F,0xB6,0xCC),
        "space_before":6})])
# shield motif on right
sh = s.shapes.add_shape(MSO_SHAPE.HEPTAGON, Inches(10.5), Inches(2.55),
                        Inches(2.0), Inches(2.4))
sh.rotation = 90
sh.fill.solid(); sh.fill.fore_color.rgb = TEAL
sh.line.fill.background(); sh.shadow.inherit = False
text(s, Inches(10.5), Inches(3.1), Inches(2.0), Inches(1.4),
     "🛡️", size=66, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# SLIDE 2 — PROBLEM & MOTIVATION
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Problem & Motivation", "Why Fraud Detection Is Hard")
footer(s, 2)
pic(s, os.path.join(ASSETS, "imbalance.png"),
    Inches(7.4), Inches(1.95), w=Inches(5.2))
bullets(s, Inches(0.6), Inches(1.9), Inches(6.5), Inches(5),
    [("Fraud is ", 0), ],
    size=16)
# rewrite first bullet via direct text block for richer formatting
text(s, Inches(0.6), Inches(1.9), Inches(6.5), Inches(4.6),
     [("Insurance fraud diverts an estimated ", {"size":16}),
      ("US$80+ billion", {"size":16,"bold":True,"color":RED}),
      (" annually worldwide — losses passed to honest policyholders as "
       "higher premiums.", {"size":16,"space_after":12})])
text(s, Inches(0.6), Inches(2.85), Inches(6.5), Inches(0.4),
     "Three structural challenges:", size=16, bold=True, color=NAVY,
     space_after=6)
bullets(s, Inches(0.6), Inches(3.35), Inches(6.5), Inches(3.2),
    [("Extreme class imbalance — only ", 0,
       {"color":INK}),
    ],
    size=15)
# replace with proper mixed bullets
text(s, Inches(0.6), Inches(3.35), Inches(6.5), Inches(3.4),
     [("▪  ", {"size":15,"bold":True,"color":TEAL}),
      ("Extreme class imbalance", {"size":15,"bold":True}),
      (" — only ", {"size":15}),
      ("6.0%", {"size":15,"bold":True,"color":RED}),
      (" of claims are fraud (923 / 15,420), so naive models score "
       "high accuracy by predicting “legitimate” for everything.",
       {"size":15,"space_after":9,"line_spacing":1.05})])
text(s, Inches(0.6), Inches(4.42), Inches(6.5), Inches(1.2),
     [("▪  ", {"size":15,"bold":True,"color":TEAL}),
      ("Black-box risk", {"size":15,"bold":True}),
      (" — regulators and investigators reject scores they cannot "
       "audit; a number alone cannot justify denying a claim.",
       {"size":15,"space_after":9,"line_spacing":1.05})])
text(s, Inches(0.6), Inches(5.30), Inches(6.5), Inches(1.2),
     [("▪  ", {"size":15,"bold":True,"color":TEAL}),
      ("Scale & latency", {"size":15,"bold":True}),
      (" — tens of thousands of daily claims must be scored in "
       "real time, then routed for investigation.",
       {"size":15,"space_after":9,"line_spacing":1.05})])


# ==================================================================
# SLIDE 3 — SYSTEM OVERVIEW / ARCHITECTURE
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "System Overview", "A Three-Tier Full-Stack Platform")
footer(s, 3)
pic(s, os.path.join(ASSETS, "pipeline.png"),
    Inches(0.5), Inches(1.65), w=Inches(12.3))
# three tier cards
tiers = [
    ("Presentation tier", "Next.js + React", TEAL,
     "Dashboards, claim intake, investigator workspace, analytics. "
     "Server-side proxy to the API."),
    ("Application tier", "FastAPI + Python", NAVY,
     "ML serving, Bayesian rule engine, three-layer explainer, "
     "LangGraph multi-agent investigation."),
    ("Data tier", "PostgreSQL + Redis", AMBER,
     "PostgreSQL persists claims, predictions, reports; Redis backs "
     "JWT sessions and rate-limiting."),
]
y0 = Inches(3.5)
for i, (t, sub, col, body) in enumerate(tiers):
    x = Inches(0.5 + i*4.25)
    card(s, x, y0, Inches(4.0), Inches(3.2), fill=LIGHT,
         line=RGBColor(0xDD,0xE3,0xEA))
    band2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x, y0, Inches(4.0), Inches(0.7))
    band2.fill.solid(); band2.fill.fore_color.rgb = col
    band2.line.fill.background(); band2.shadow.inherit = False
    try: band2.adjustments[0] = 0.08
    except Exception: pass
    text(s, x, y0+Emu(60000), Inches(4.0), Inches(0.7), t, size=15,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.2), y0+Inches(0.9), Inches(3.6), Inches(0.4),
         sub, size=13, bold=True, color=col)
    text(s, x+Inches(0.2), y0+Inches(1.35), Inches(3.6), Inches(1.8),
         body, size=12, color=INK, line_spacing=1.1)
text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
     "Docker Compose orchestrates all four services; a single "
     "`docker compose up -d` brings the full stack online.",
     size=11, color=SLATE, align=PP_ALIGN.CENTER)


# ==================================================================
# SLIDE 4 — DATA & FEATURES
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Data & Feature Engineering", "From 31 Raw Columns to 116 Model Inputs")
footer(s, 4)
# KPI row
kpi(s, Inches(0.5), Inches(1.7), Inches(2.8), Inches(1.4),
    "15,420", "claims (1,542 hold-out)", vcolor=NAVY)
kpi(s, Inches(3.55), Inches(1.7), Inches(2.8), Inches(1.4),
    "116", "model features", vcolor=TEAL)
kpi(s, Inches(6.6), Inches(1.7), Inches(2.8), Inches(1.4),
    "6", "engineered features", vcolor=AMBER)
kpi(s, Inches(9.65), Inches(1.7), Inches(3.2), Inches(1.4),
    "ROC-AUC 0.83", "best baseline model", vcolor=GREEN)
# left: raw features
card(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(3.5),
     fill=LIGHT, line=RGBColor(0xDD,0xE3,0xEA))
text(s, Inches(0.7), Inches(3.55), Inches(5.6), Inches(0.4),
     "Raw attributes (24 categorical · 7 numeric)", size=13,
     bold=True, color=NAVY)
text(s, Inches(0.7), Inches(4.05), Inches(5.6), Inches(2.8),
     [("Driver — Age, Sex, MaritalStatus, DriverRating, "
       "PastNumberOfClaims", {"size":11.5,"space_after":5}),
      ("Policy — BasePolicy, PolicyType, Deductible, "
       "Days:Policy-Accident", {"size":11.5,"space_after":5}),
      ("Vehicle — Make, Category, Price, AgeOfVehicle", 
       {"size":11.5,"space_after":5}),
      ("Accident — Month, DayOfWeek, Area, Fault, "
       "NumberOfCars", {"size":11.5,"space_after":5}),
      ("Claim — PoliceReportFiled, AgentType, "
       "AddressChange-Claim, NumberOfSuppliments",
       {"size":11.5,"space_after":5})],
     size=11.5, color=INK, line_spacing=1.1)
# right: engineered
card(s, Inches(6.8), Inches(3.4), Inches(6.0), Inches(3.5),
     fill=LIGHT, line=RGBColor(0xDD,0xE3,0xEA))
text(s, Inches(7.0), Inches(3.55), Inches(5.6), Inches(0.4),
     "Domain-engineered features", size=13, bold=True, color=TEAL)
text(s, Inches(7.0), Inches(4.05), Inches(5.6), Inches(2.8),
     [("vehicle_age_engineered", {"size":12,"bold":True}),
      (" — proxy for vehicle depreciation", {"size":11.5,"space_after":5}),
      ("policy_accident_delay_engineered", {"size":12,"bold":True}),
      (" — early-accident-on-new-policy signal", {"size":11.5,"space_after":5}),
      ("claim_report_delay_engineered", {"size":12,"bold":True}),
      (" — late reporting risk", {"size":11.5,"space_after":5}),
      ("deductible_vehicle_ratio_engineered", {"size":12,"bold":True}),
      (" — high-deductible / low-value anomaly",
       {"size":11.5,"space_after":5}),
      ("weekend_accident_engineered", {"size":12,"bold":True}),
      ("  ·  ", {"size":11.5}),
      ("high_value_single_vehicle_engineered",
       {"size":12,"bold":True})],
     size=11.5, color=INK, line_spacing=1.1)


# ==================================================================
# SLIDE 5 — MODELING PIPELINE
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Modeling Pipeline", "Two-Stage Tuning with Pareto-Front Selection")
footer(s, 5)
pic(s, os.path.join(ASSETS, "leaderboard.png"),
    Inches(0.5), Inches(1.7), w=Inches(7.6))
# right column narrative
text(s, Inches(8.4), Inches(1.75), Inches(4.5), Inches(0.4),
     "Method", size=14, bold=True, color=NAVY)
text(s, Inches(8.4), Inches(2.2), Inches(4.5), Inches(4.5),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("6 algorithms × 2 samplers", {"size":12.5,"bold":True}),
      (" baseline sweep (LightGBM, XGBoost, GBM under SMOTE "
       "and CTGAN).", {"size":12.5,"space_after":7,"line_spacing":1.05})])
text(s, Inches(8.4), Inches(3.35), Inches(4.5), Inches(2),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("Hyper-parameter tuning + ", {"size":12.5}),
      ("probability calibration", {"size":12.5,"bold":True}),
      (" on the top-3, evaluated on a held-out test set.",
       {"size":12.5,"space_after":7,"line_spacing":1.05})])
text(s, Inches(8.4), Inches(4.65), Inches(4.5), Inches(2),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("Pareto front", {"size":12.5,"bold":True}),
      (" selection across Precision, Recall, FPR and latency; "
       "production model chosen by composite score.",
       {"size":12.5,"space_after":7,"line_spacing":1.05})])
text(s, Inches(0.5), Inches(6.55), Inches(7.6), Inches(0.5),
     [("Production model: ", {"size":12,"bold":True,"color":NAVY}),
      ("CTGAN + LightGBM", {"size":12,"bold":True,"color":TEAL}),
      ("  ·  Recall 0.209  ·  FPR 0.029  ·  Brier 0.054  ·  "
       "~243 ms latency", {"size":12,"color":INK})])


# ==================================================================
# SLIDE 6 — RULE ENGINE (hybrid)
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Hybrid Decision Layer", "ML Probability + Bayesian Rule Engine")
footer(s, 6)
pic(s, os.path.join(ASSETS, "tiers.png"),
    Inches(0.5), Inches(1.55), w=Inches(12.3))
text(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.5),
     "Each triggered rule reduces the remaining “innocence” — "
     "scores compound with diminishing returns, never balloon.",
     size=12.5, color=SLATE, align=PP_ALIGN.CENTER)
# equation card
card(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.8),
     fill=NAVY, line=NAVY)
text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.7),
     [("P(final) = 1 − (1 − P(ML)) · ∏ (1 − scoreᵢ)",
       {"size":17,"bold":True,"color":WHITE})],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# rule table-ish
rules = [
    ("INVALID_POLICY / VEHICLE", "0.35", RED),
    ("JUNK / SUSPICIOUS DOCUMENT", "0.15–0.20", AMBER),
    ("FIR CLAIMED, NO NUMBER", "0.15", AMBER),
    ("NO POLICE FIR", "0.12", AMBER),
    ("SUSPICIOUS / WRONG DATA", "0.10", SLATE),
    ("MISSING DATA (per field)", "0.08", SLATE),
    ("INJURED, NO MEDICAL", "0.08", SLATE),
    ("NO WITNESS", "0.05", SLATE),
]
x0 = Inches(0.5); y0 = Inches(4.25)
cw = Inches(3.0); ch = Inches(1.25); gx = Inches(0.12); gy = Inches(0.15)
for i, (name, sc, col) in enumerate(rules):
    r, c = divmod(i, 4)
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    card(s, x, y, cw, ch, fill=LIGHT, line=RGBColor(0xDD,0xE3,0xEA))
    bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), ch)
    bb.fill.solid(); bb.fill.fore_color.rgb = col
    bb.line.fill.background(); bb.shadow.inherit = False
    text(s, x+Inches(0.25), y+Inches(0.12), cw-Inches(0.35), Inches(0.7),
         name, size=10.5, bold=True, color=NAVY, line_spacing=0.95)
    text(s, x+Inches(0.25), y+Inches(0.78), cw-Inches(0.35), Inches(0.4),
         "score_added = " + sc, size=10, color=col, bold=True)


# ==================================================================
# SLIDE 7 — EXPLAINABILITY
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Explainable AI", "Three Layers Turn a Score into an Argument")
footer(s, 7)
pic(s, os.path.join(ASSETS, "explainability.png"),
    Inches(0.6), Inches(1.7), w=Inches(8.2))
# example callout
card(s, Inches(9.1), Inches(1.7), Inches(3.8), Inches(5.0),
     fill=NAVY, line=NAVY)
text(s, Inches(9.3), Inches(1.85), Inches(3.4), Inches(0.4),
     "EXAMPLE OUTPUT", size=10, bold=True, color=AMBER)
text(s, Inches(9.3), Inches(2.25), Inches(3.4), Inches(4.3),
     [("Fraud probability ", {"size":12,"color":RGBColor(0xCF,0xDD,0xEA)}),
      ("94%", {"size":20,"bold":True,"color":RED}),
      ("  —  Critical", {"size":12,"bold":True,"color":RED,
        "space_after":10}),
      ("“The policyholder was at fault, no police report was "
       "filed, and the policy number is not in our records. "
       "A high deductible-to-vehicle ratio further raises "
       "risk.”", {"size":11.5,"color":WHITE,"line_spacing":1.15,
        "space_after":10,"italic":False}),
      ("SHAP top driver:  Deductible  (+1.4)",
       {"size":10.5,"color":AMBER,"bold":True,"space_after":4}),
      ("Rule flags:  unregistered policy · no FIR",
       {"size":10.5,"color":AMBER,"bold":True})],
     line_spacing=1.1)
text(s, Inches(0.6), Inches(6.5), Inches(8.2), Inches(0.4),
     "Layer 1 is deterministic & reproducible; Layer 2 (LLM) is "
     "optional with a deterministic fallback.", size=10.5,
     color=SLATE)


# ==================================================================
# SLIDE 8 — MULTI-AGENT
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Multi-Agent Investigation", "12 LangGraph Agents, Parallel Fan-Out")
footer(s, 8)
pic(s, os.path.join(ASSETS, "agents.png"),
    Inches(0.4), Inches(1.55), w=Inches(8.7))
# right side: phases
text(s, Inches(9.3), Inches(1.7), Inches(3.6), Inches(0.4),
     "Pipeline phases", size=14, bold=True, color=NAVY)
phases = [
    ("① Fan-out", "6 data-gathering agents run concurrently "
     "(claim, prediction, SHAP, knowledge, policy, history).", TEAL),
    ("② Synthesis", "Contradiction, Risk, Recommendation agents "
     "cross-check the evidence.", AMBER),
    ("③ Reflection & Report", "Reflection self-critique → draft "
     "report → QA agent validates.", NAVY),
    ("④ Human-in-the-loop", "Graph interrupts before storage; an "
     "investigator approves or edits.", GREEN),
]
y = Inches(2.2)
for t, body, col in phases:
    text(s, Inches(9.3), y, Inches(3.7), Inches(0.35),
         t, size=12, bold=True, color=col)
    text(s, Inches(9.3), y+Inches(0.33), Inches(3.7), Inches(0.85),
         body, size=10.5, color=INK, line_spacing=1.05)
    y += Inches(1.18)


# ==================================================================
# SLIDE 9 — EVALUATION / RESULTS
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Evaluation", "Operating Trade-offs Across Samplers")
footer(s, 9)
pic(s, os.path.join(ASSETS, "tradeoff.png"),
    Inches(0.5), Inches(1.7), w=Inches(7.8))
# KPI strip
metrics = [
    ("0.83", "ROC-AUC", NAVY),
    ("0.054", "Brier score", TEAL),
    ("87 ms", "Fastest model", AMBER),
    ("×5", "Fraud recall vs\nmajority baseline", GREEN),
]
x = Inches(0.6)
for v, l, col in metrics:
    kpi(s, x, Inches(5.7), Inches(1.8), Inches(1.25), v, l, vcolor=col)
    x += Inches(1.95)
# takeaways
card(s, Inches(8.55), Inches(1.7), Inches(4.3), Inches(5.25),
     fill=LIGHT, line=RGBColor(0xDD,0xE3,0xEA))
text(s, Inches(8.75), Inches(1.85), Inches(3.9), Inches(0.4),
     "Findings", size=14, bold=True, color=NAVY)
text(s, Inches(8.75), Inches(2.35), Inches(3.9), Inches(4.5),
     [("▪  ", {"size":12,"bold":True,"color":TEAL}),
      ("CTGAN synthesizing minority examples lifted ", {"size":12}),
      ("recall", {"size":12,"bold":True}),
      (" versus SMOTE at comparable precision.",
       {"size":12,"space_after":9,"line_spacing":1.1})])
text(s, Inches(8.75), Inches(3.65), Inches(3.9), Inches(2),
     [("▪  ", {"size":12,"bold":True,"color":TEAL}),
      ("SMOTE+LightGBM gave the ", {"size":12}),
      ("lowest FPR", {"size":12,"bold":True}),
      (" (0.015) — ideal when false alarms are costly.",
       {"size":12,"space_after":9,"line_spacing":1.1})])
text(s, Inches(8.75), Inches(5.0), Inches(3.9), Inches(2),
     [("▪  ", {"size":12,"bold":True,"color":TEAL}),
      ("Calibration (Brier 0.054) makes probabilities "
       "trustworthy for tiering and investigation routing.",
       {"size":12,"space_after":9,"line_spacing":1.1})])


# ==================================================================
# SLIDE 10 — DEPLOYMENT & ENGINEERING
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
header(s, "Deployment & Engineering", "Production-Grade Operations")
footer(s, 10)
cols = [
    ("Containerised stack", "🐳", NAVY,
     ["Single `docker compose up` brings up PostgreSQL, Redis, "
      "FastAPI backend and Next.js frontend.",
      "Health checks and dependency ordering on all services."]),
    ("MLOps lifecycle", "🔄", TEAL,
     ["Trained pipeline serialized as `.pkl`; `features.json` is the "
      "authoritative schema.",
      "Singleton FraudModelService loads the Rank-1 tuned checkpoint "
      "and a model registry version."]),
    ("Security & auth", "🔐", AMBER,
     ["JWT auth with Redis-backed session store and rate limiting.",
      "Role-based access; default admin rotated on deployment."]),
    ("Observability", "📊", GREEN,
     ["Drift detection: compares 7-day High/Critical rate to a 15% "
      "baseline.",
      "Audit logging of predictions and agent runs for regulators."]),
]
x = Inches(0.5)
for t, icon, col, items in cols:
    card(s, x, Inches(1.75), Inches(3.0), Inches(5.0),
         fill=LIGHT, line=RGBColor(0xDD,0xE3,0xEA))
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, Inches(1.75), Inches(3.0), Inches(0.7))
    cb.fill.solid(); cb.fill.fore_color.rgb = col
    cb.line.fill.background(); cb.shadow.inherit = False
    try: cb.adjustments[0] = 0.08
    except Exception: pass
    text(s, x, Inches(1.8), Inches(3.0), Inches(0.6),
         icon + "  " + t, size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = Inches(2.65)
    for it in items:
        text(s, x+Inches(0.18), yy, Inches(2.7), Inches(2),
             [("▪  ", {"size":11,"bold":True,"color":col}),
              (it, {"size":11,"color":INK,"line_spacing":1.08})],
             line_spacing=1.08)
        yy += Inches(1.95)
    x += Inches(3.15)


# ==================================================================
# SLIDE 11 — CONCLUSION & FUTURE WORK
# ==================================================================
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
# accent
ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1),
                        Inches(1.2), Pt(5))
ab.fill.solid(); ab.fill.fore_color.rgb = AMBER
ab.line.fill.background(); ab.shadow.inherit = False
text(s, Inches(0.8), Inches(0.55), Inches(11), Inches(0.4),
     "CONCLUSION", size=12, bold=True, color=TEAL)
text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1.0),
     "A defensible path from raw claim to investigated decision",
     size=30, bold=True, color=WHITE)
# contributions
card(s, Inches(0.8), Inches(2.7), Inches(5.8), Inches(3.4),
     fill=RGBColor(0x10,0x2A,0x4D), line=TEAL, line_w=1.0)
text(s, Inches(1.0), Inches(2.85), Inches(5.4), Inches(0.4),
     "Contributions", size=15, bold=True, color=AMBER)
text(s, Inches(1.0), Inches(3.35), Inches(5.4), Inches(2.7),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("Hybrid scoring — calibrated LightGBM fused with a "
       "Bayesian, auditable rule engine.",
       {"size":12.5,"color":WHITE,"space_after":9,"line_spacing":1.12})])
text(s, Inches(1.0), Inches(4.45), Inches(5.4), Inches(1.5),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("Three-layer explainability — SHAP + humaniser + LLM "
       "narrative with deterministic fallback.",
       {"size":12.5,"color":WHITE,"space_after":9,"line_spacing":1.12})])
text(s, Inches(1.0), Inches(5.55), Inches(5.4), Inches(1.0),
     [("▪  ", {"size":12.5,"bold":True,"color":TEAL}),
      ("12-agent LangGraph investigation with human-in-the-loop "
       "approval.", {"size":12.5,"color":WHITE,"line_spacing":1.12})])
# future work
card(s, Inches(6.9), Inches(2.7), Inches(5.8), Inches(3.4),
     fill=RGBColor(0x10,0x2A,0x4D), line=AMBER, line_w=1.0)
text(s, Inches(7.1), Inches(2.85), Inches(5.4), Inches(0.4),
     "Future work", size=15, bold=True, color=AMBER)
text(s, Inches(7.1), Inches(3.35), Inches(5.4), Inches(2.7),
     [("▪  ", {"size":12.5,"bold":True,"color":AMBER}),
      ("Online learning & drift-triggered retraining from "
       "investigator feedback.",
       {"size":12.5,"color":WHITE,"space_after":9,"line_spacing":1.12})])
text(s, Inches(7.1), Inches(4.4), Inches(5.4), Inches(1.5),
     [("▪  ", {"size":12.5,"bold":True,"color":AMBER}),
      ("Graph-based claim-network features for organised-fraud "
       "rings.", {"size":12.5,"color":WHITE,"space_after":9,
        "line_spacing":1.12})])
text(s, Inches(7.1), Inches(5.45), Inches(5.4), Inches(1.0),
     [("▪  ", {"size":12.5,"bold":True,"color":AMBER}),
      ("Adversarial-robustness and bias audits across demographic "
       "groups.", {"size":12.5,"color":WHITE,"line_spacing":1.12})])
text(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5),
     "Thank you — questions welcome.", size=16, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("Saved", OUT)
print("Slides:", len(prs.slides._sldIdLst))
